import os
from pathlib import Path
import docx
from pptx import Presentation
from markitdown import MarkItDown

def extract_images_from_docx(docx_path: Path, media_dir: Path) -> int:
    """Extracts embedded images from a Word Document (.docx)."""
    doc = docx.Document(docx_path)
    image_count = 0
    for rel_id, rel in doc.part.rels.items():
        if "image" in rel.target_ref:
            image_bytes = rel.target_part.blob
            image_count += 1
            image_ext = Path(rel.target_ref).suffix or ".png"
            with open(media_dir / f"image_{image_count}{image_ext}", "wb") as img_f:
                img_f.write(image_bytes)
    return image_count

def extract_images_from_pptx(pptx_path: Path, media_dir: Path) -> int:
    """Extracts embedded images from a PowerPoint Presentation (.pptx)."""
    prs = Presentation(pptx_path)
    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 represents a Picture Shape
                image = shape.image
                image_bytes = image.blob
                image_count += 1
                image_ext = image.ext or "png"
                with open(media_dir / f"image_{image_count}.{image_ext}", "wb") as img_f:
                    img_f.write(image_bytes)
    return image_count

def batch_convert_specifications(input_dir: str, output_dir: str):
    """
    Scans input_dir for .docx, .xlsx, and .pptx specifications.
    Converts them to clean Markdown inside unique, isolated folders.
    """
    md_engine = MarkItDown()
    input_path = Path(input_dir)
    output_base_path = Path(output_dir)
    
    supported_extensions = [".docx", ".xlsx", ".xls", ".pptx"]
    files_to_process = [f for f in input_path.iterdir() if f.suffix.lower() in supported_extensions]
    
    if not files_to_process:
        print(f"📁 No matching files found in '{input_dir}'.")
        return

    print(f"🚀 Found {len(files_to_process)} technical specifications. Starting multi-format pipeline...")
    print("-" * 75)

    for file_path in files_to_process:
        doc_name = file_path.stem
        extension = file_path.suffix.lower()
        
        # 1. Establish an isolated sandbox directory per document (Always created)
        unique_doc_dir = output_base_path / doc_name
        unique_doc_dir.mkdir(parents=True, exist_ok=True)
        
        media_dir = unique_doc_dir / "media"
        target_md_path = unique_doc_dir / f"{doc_name}.md"
        
        print(f"⏳ Processing [{extension.upper()}]: {file_path.name}")
        
        try:
            # 2. Extract media based on file type
            extracted_images = 0
            if extension == ".docx":
                media_dir.mkdir(parents=True, exist_ok=True)
                extracted_images = extract_images_from_docx(file_path, media_dir)
            elif extension == ".pptx":
                media_dir.mkdir(parents=True, exist_ok=True)
                extracted_images = extract_images_from_pptx(file_path, media_dir)
            
            # 3. Parse core text layout to Markdown
            result = md_engine.convert(str(file_path))
            markdown_content = result.text_content
            
            # 4. Clean messy Base64 chunks and replace with structural image placeholders
            cleaned_lines = []
            img_index = 1
            for line in markdown_content.splitlines():
                if "data:image" in line and "base64" in line:
                    # Make sure media folder exists if an excel sheet contains an embedded chart/image asset
                    media_dir.mkdir(parents=True, exist_ok=True)
                    cleaned_lines.append(
                        f"\n\n![Specification Image {img_index}](media/image_{img_index}.png)\n"
                        f"> **[Copilot Context Placeholder]**: Provide descriptions for `media/image_{img_index}.png` here to ensure complete test generation coverage.\n"
                    )
                    img_index += 1
                else:
                    cleaned_lines.append(line)
            
            # Save the clean Markdown file
            with open(target_md_path, "w", encoding="utf-8") as f:
                f.write("\n".join(cleaned_lines))
                
            print(f"✅ Folder isolated: {unique_doc_dir.name}")
            if extracted_images > 0 or img_index > 1:
                print(f"   ↳ Clean Markdown saved. Image assets linked to /media.")
            else:
                print(f"   ↳ Clean Markdown data table/text parsed successfully.")
            print("-" * 75)
            
        except Exception as e:
            print(f"❌ Failed to process {file_path.name}: {str(e)}")
            print("-" * 75)
            
    print("✨ Multi-format pipeline processing complete!")

if __name__ == "__main__":
    INPUT_FOLDER = "docs/raw_specs"
    OUTPUT_FOLDER = "docs/specifications"
    
    batch_convert_specifications(INPUT_FOLDER, OUTPUT_FOLDER)